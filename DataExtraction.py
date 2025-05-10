
from pptx import Presentation as PptxPresentation
import uuid
from spire.presentation import *


class SlideObj:
    def __init__(self, text, slide):
        self.text=text
        self.slide=slide

    def get_text(self):
        return str(self.text)

slides={}
#
def extract_text_from_pptx(pptx_path):
    prs = PptxPresentation(pptx_path)
    sourcePPT = Presentation()
    sourcePPT.LoadFromFile(pptx_path)
    index=0;

    for slide in prs.slides:
        all_text = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        all_text+=run.text+" "


        id=str(uuid.uuid4())
        text=all_text.replace('\t',"").replace("  "," ")
        curr_slide=SlideObj(text, sourcePPT.Slides[index])
        slides[id]=curr_slide
        index+=1



extract_text_from_pptx("data/inputPPTs/Advantages-and-Disadvantages-of-AI.pptx")
extract_text_from_pptx("data/inputPPTs/AI.pptx.pptx")
extract_text_from_pptx("data/inputPPTs/Cyber-Security-Awarness-Slide-September-2022 (1).pptx")
extract_text_from_pptx("data/inputPPTs/cybersecurity.pptx")
extract_text_from_pptx("data/inputPPTs/IT-Security-20210426203847.pptx")




