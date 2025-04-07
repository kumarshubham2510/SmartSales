from CreateVectorData import results
from main import slides
from pptx import Presentation

prs = Presentation()

slides_to_copy=[slides[x].slide for x in results]

def add_slides_from_list(slide_objects):

    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[5]  # Choose a blank layout

    for slide in slide_objects:
        # Add a new blank slide to the presentation
        new_slide = prs.slides.add_slide(blank_slide_layout)
        #Clear Add to Title textbox from blank slides
        for shape in list(new_slide.shapes):
            new_slide.shapes._spTree.remove(shape._element)

        # Copy each shape from the source slide to the new slide
        for shape in slide.shapes:
            el = shape.element
            newel = new_slide.shapes.element
            newel.append(el)

        # Copy slide notes if they exist
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
            new_notes_slide = new_slide.notes_slide
            new_notes_slide.notes_text_frame.text = notes

    prs.save("./GeneratedResults/new_presentation.pptx")
    print("New presentation 'new_presentation.pptx' created with added slides in the GeneratedResults ppt")

add_slides_from_list(slides_to_copy)





