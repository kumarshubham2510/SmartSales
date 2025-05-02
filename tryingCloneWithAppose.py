
import aspose.slides as slides
from pptx import Presentation

srcPres = "data/inputPPTs/AI.pptx.pptx"
# OUTPUT_PATH= "new_presentation.pptx"
#
# prs = Presentation()
#
# blank_slide_layout = prs.slide_layouts[6]
#
# # Add a blank slide
# slide = prs.slides.add_slide(blank_slide_layout)
#
# # Save the presentation
# prs.save(OUTPUT_PATH)

# Instantiate Presentation class to load the source presentation file
with slides.Presentation("data/inputPPTs/AI.pptx.pptx") as srcPres:
    # Instantiate Presentation class for destination PPTX (where slide is to be cloned)
    with slides.Presentation() as destPres:
        # Clone the desired slide from the source presentation to the end of the collection of slides in destination presentation
        slds = destPres.slides
        slds.add_clone(srcPres.slides[4])

        # Write the destination presentation to disk
        destPres.save("Aspose2_out.pptx", slides.export.SaveFormat.PPTX)