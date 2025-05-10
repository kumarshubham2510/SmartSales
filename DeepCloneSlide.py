from pptx import Presentation as PptxPresetation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from io import BytesIO
from spire.presentation import *


def add_slides_from_list(slide_objects):
    destPPT = Presentation()
    for slide in slide_objects:
        destPPT.Slides.AppendBySlide(slide)


    return destPPT










