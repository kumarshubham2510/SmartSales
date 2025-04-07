
from pptx import Presentation
import uuid


class SlideObj:
    def __init__(self, text, slide):
        self.text=text
        self.slide=slide

    def get_text(self):
        return str(self.text)

slides={}
#
def extract_text_from_pptx(pptx_path):
    prs = Presentation(pptx_path)

    for slide in prs.slides:
        all_text = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        all_text+=run.text+" "

        id=str(uuid.uuid4())
        text=all_text.replace('\t',"").replace("  "," ")
        curr_slide=SlideObj(text, slide)
        slides[id]=curr_slide



extract_text_from_pptx("./PPTs/Advantages-and-Disadvantages-of-AI.pptx")
extract_text_from_pptx("./PPTs/AI.pptx.pptx")
extract_text_from_pptx("./PPTs/Cyber-Security-Awarness-Slide-September-2022 (1).pptx")
extract_text_from_pptx("./PPTs/cybersecurity.pptx")
extract_text_from_pptx("./PPTs/IT-Security-20210426203847.pptx")




